import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- COLOR PALETTE ---
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
BLUE = RGBColor(0x2E, 0x75, 0xB6)
GREEN = RGBColor(0x70, 0xAD, 0x47)
AMBER = RGBColor(0xFF, 0xC0, 0x00)
TEAL = RGBColor(0x1C, 0x72, 0x93)
GREY = RGBColor(0x80, 0x80, 0x80)


def set_text_format(run, font_size, bold=False, color=NAVY):
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Segoe UI'


def add_image_placeholder(slide, left, top, width, height, text):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)  # Light grey
    shape.line.color.rgb = GREY

    tf = shape.text_frame
    tf.text = f"IMAGE PLACEHOLDER:\n{text}"
    for paragraph in tf.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            set_text_format(run, 12, bold=True, color=GREY)


def create_deck():
    prs = Presentation()
    # 16:9 Aspect Ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    layout_title = prs.slide_layouts[0]
    layout_content = prs.slide_layouts[1]
    layout_blank = prs.slide_layouts[6]

    # ==========================================
    # SLIDE 1: TITLE
    # ==========================================
    slide1 = prs.slides.add_slide(layout_blank)

    # Title
    txBox = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(11.3), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Mechatronics Modelling of an Entire Industry 5.0 Production Line"
    set_text_format(p.runs[0], 40, bold=True, color=NAVY)

    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = "A validated multiphysics digital model — Modelica / Dymola · FMI 2.0 · Python"
    set_text_format(p2.runs[0], 24, color=TEAL)

    # Presenter Details
    txBox2 = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(10), Inches(2))
    tf2 = txBox2.text_frame
    details = [
        "Samiksha Satre | N. K. Orchid College of Engineering & Technology, Solapur",
        "Euler Laboratory, ISAE-Supméca | Supervisor: Mr. Romain Delabeye",
        "February–June 2026",
        "",
        "Prepared for the Dassault Systèmes review meeting."
    ]
    for i, detail in enumerate(details):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = detail
        if len(p.runs) > 0:
            set_text_format(p.runs[0], 16, bold=(i == 0), color=NAVY)
    # Placeholders
    add_image_placeholder(slide1, 0.5, 0.5, 2, 1, "ISAE-Supméca Logo")
    add_image_placeholder(slide1, 10.8, 0.5, 2, 1, "College Logo")
    add_image_placeholder(slide1, 10.8, 6, 2, 1, "Dassault Systèmes Logo")

    # ==========================================
    # HELPER FOR STANDARD SLIDES (2, 3, 5, 6, 7, 8, 9, 10, 11, 13, 14, 15, 17, 18)
    # ==========================================
    def add_standard_slide(title_text, body_text_list, placeholder_args=None):
        slide = prs.slides.add_slide(layout_content)
        title = slide.shapes.title
        title.text = title_text
        set_text_format(title.text_frame.paragraphs[0].runs[0], 32, bold=True, color=NAVY)

        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.text = body_text_list[0]
        set_text_format(tf.paragraphs[0].runs[0], 18, color=NAVY)

        for point in body_text_list[1:]:
            p = tf.add_paragraph()
            p.text = point
            if not p.runs:
                run = p.add_run()
                run.text = p.text
                p.text = ""
            else:
                run = p.runs[0]

            set_text_format(run, 18, color=NAVY)
            p.level = 0 if not point.startswith("  ") else 1

        if placeholder_args:
            add_image_placeholder(slide, *placeholder_args)
        return slide

    # SLIDE 2
    add_standard_slide("THE MISSION: Bridging Physical and Digital", [
        "OBJECTIVE:",
        "  Build a unified process, mechatronic, and multiphysics model of an entire",
        "  Industry 5.0 production line to support smart reconfiguration and optimization.",
        "",
        "WORKFLOW:",
        "1. Real-time experiment on ALIX line",
        "2. TDMS data acquisition",
        "3. Multiphysics modelling in Dymola",
        "4. FMU export (FMI 2.0 Co-Simulation)",
        "5. Python orchestration",
        "6. Parameter identification & validation",
        "",
        "HOST: Euler Laboratory, ISAE-Supméca · IS2M theme"
    ])

    # SLIDE 3
    add_standard_slide("THE SYSTEM: ALIX / ERMASMART Line", [
        "Specifications: 8445 mm length · 3-phase 400 V / 50 Hz · S7-1200+TULIP · 30 s cadence",
        "",
        "THE FIVE STATIONS:",
        "• ON10 (Bin-pick) — Out of scope",
        "• DX10 (Dosing) — Out of scope",
        "• MI00+MD20 (CR5 capping cobot) — MODELLED (Multibody)",
        "• XY10 (Cartesian XYZ pick-and-place) — MODELLED (Cartesian)",
        "• VL10 (Vertical storage) — Out of scope",
        "",
        "CONTINUOUS CONVEYOR BELT — MODELLED (Multiphysics)"
    ], (9, 2.5, 4, 4, "alix_spatial_layout.png thumbnail"))

    # ==========================================
    # SLIDE 4: SCOPE (TABLE)
    # ==========================================
    slide4 = prs.slides.add_slide(layout_blank)
    title = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(10), Inches(1))
    title.text_frame.text = "SCOPE: Defined by Experimental Validation Capability"
    set_text_format(title.text_frame.paragraphs[0].runs[0], 32, bold=True, color=NAVY)

    rows, cols = 7, 5
    table_shape = slide4.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(12), Inches(4))
    table = table_shape.table

    headers = ["Station", "Hands-on Access", "TDMS Data", "Documentation", "Modelled"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h

    data = [
        ["Conveyor", "Yes", "Yes", "Yes", "YES (Blue)"],
        ["CR5 Cobot", "Yes", "Yes", "Yes", "YES (Green)"],
        ["XY10 Station", "Yes", "Yes", "Yes", "YES (Amber)"],
        ["ON10 Bin-pick", "No", "No", "No", "No (Grey)"],
        ["DX10 Dosing", "No", "No", "No", "No (Grey)"],
        ["VL10 Storage", "No", "No", "No", "No (Grey)"]
    ]

    for r in range(6):
        for c in range(5):
            table.cell(r + 1, c).text = data[r][c]

    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(6), Inches(10), Inches(1))
    txBox.text_frame.text = "Principle: A model that cannot be validated is a guess.\nSupervisor: 'Scope is perfect! Evaluate performance metrics.' — Romain Delabeye"

    # SLIDE 5
    add_standard_slide("THE METHOD: 6-Stage Workflow", [
        "1. Real-time experiment on ALIX: XY10 picks → conveyor transports → CR5 caps.",
        "2. TDMS data acquisition: 6-channel NI DAQ at 6 kHz, U×200 / I×10 gains.",
        "3. Multiphysics modelling in Dymola on 3DEXPERIENCE. Parameters traceable.",
        "4. FMU export: FMI 2.0 Co-Simulation, solver DASSL, tolerance 10⁻⁴.",
        "5. Python orchestration: master.py + fmu_manager.py using FMpy library.",
        "6. Two-stage validation: Stage A (FMU vs TDMS) and Stage B (FMU vs Dymola).",
        "",
        "Built entirely on Dymola / 3DEXPERIENCE · exported as FMI 2.0 Co-Simulation FMUs",
        "Orchestrated from Python — your stack, end-to-end."
    ])

    # SLIDE 6
    add_standard_slide("THE MODELS: Digital Twin Library Breakdown", [
        "1. CONVEYOR (Full Multiphysics)",
        "  • Electrical + Mechanical + Thermal. 15 blocks, 161 parameters.",
        "",
        "2. CR5 COBOT (6-DOF Multibody)",
        "  • Newton-Euler equations auto-assembled. Mass/inertia from CAD.",
        "",
        "3. XY10 STATION (Cartesian + Vacuum)",
        "  • Mechanical + Electrical + Pneumatic. DC-equivalent drives.",
        "",
        "Note: Conveyor is continuous; CR5 and XY10 act in bursts."
    ])

    # SLIDE 7
    add_standard_slide("CONVEYOR MODEL DETAIL", [
        "ARCHITECTURE:",
        "Supply 244 V → Induction motor → Gearbox i=20 → Drum → Belt",
        "Thermal Branch: Motor heat capacity (500 J/K) → Conductance (1.5 W/K) → Ambient 20 °C",
        "",
        "PARAMETERS:",
        "Rs=12 Ω | Rr=8 Ω | Lm=0.45 H | Lσ=0.025 H | ratio=20 | C=500 J/K | G=1.5 W/K | τ=333 s",
        "",
        "STEADY-STATE OUTPUT:",
        "1499 rpm | 1.629 A RMS | 107.96 W | +11.95 K rise"
    ], (8, 2, 5, 4, "Dymola Diagram Conveyor_updated_v3.mo"))

    # SLIDE 8
    add_standard_slide("CR5 COBOT MODEL DETAIL", [
        "KINEMATIC CHAIN:",
        "Base → J1(Z) → Link1 → J2(Y) → Link2 → J3(Y) → Link3 → J4(X) → Link4 → J5(Y) → Link5 → J6(Z) → End-effector",
        "",
        "GEOMETRY EXTRACTED FROM REAL ROBOT CAD:",
        "• Mass, COM, inertia per link extracted via 3DEXPERIENCE Measure Inertia.",
        "• Modelica MultiBody joints automatically assemble Newton-Euler equations.",
        "",
        "PHYSICAL CHECKS:",
        "• Total mass agreement: 23.74 kg vs brochure 25 kg = 94.9 %",
        "• Joint-2 gravity pendulum test passes perfectly."
    ], (8, 2, 5, 4, "CR5 CAD render / cr5_J2_pendulum.png"))

    # SLIDE 9
    add_standard_slide("XY10 MODEL DETAIL", [
        "CLOSED-LOOP CHAIN (X, Y, Z):",
        "x_cmd → PID → Saturation ±V_dc → SignalVoltage → R-L Winding → CurrentSensor → IdealGearR2T → Mass",
        "",
        "HONEST ABSTRACTION: DC-EQUIVALENT STEPPERS",
        "• Real chopper produces ±20 A; FMU produces smooth time-averaged ±3 A.",
        "• Envelope, timing and energy are captured perfectly.",
        "• Chopper detail belongs to EMC studies, not digital-twin operation.",
        "",
        "SANITY CHECK:",
        "• Zero-command equilibrium: Z droops to −5.2 mm with 20.8 mA holding current (gravity verified)."
    ], (8, 2, 5, 4, "Dymola Diagram XY10Station_Diagram.mo"))

    # SLIDE 10
    add_standard_slide("VALIDATION DATA (TDMS)", [
        "PHYSICAL EXPERIMENT:",
        "1. XY10 picks pot (vacuum)",
        "2. XY10 places on belt",
        "3. Conveyor transports",
        "4. CR5 caps pot",
        "5. TDMS records U0-U2 & I0-I2",
        "",
        "MEASUREMENT SETUP:",
        "• 6-channel NI DAQ at 6 kHz",
        "• Transducer gains: ×200 voltage / ×10 current",
        "• Minimal-instrumentation principle: only the line's electrical signature used."
    ], (8, 2, 5, 4, "Lab bench photo / TDMS icon"))

    # SLIDE 11
    add_standard_slide("VALIDATION STRATEGY", [
        "STAGE A: MODEL VS WORLD",
        "• FMU vs TDMS: Does the model match physical reality?",
        "",
        "STAGE B: EXPORT FAITHFULNESS",
        "• FMU vs Dymola: Did the export preserve the math without loss?",
        "",
        "7-STEP PIPELINE:",
        "Load TDMS → Apply gains → Capture window → Band-pass 40–60 Hz → Edge-crop →",
        "Decimate FMU 50 kHz→6 kHz → Cross-correlate phase-align",
        "",
        "Metrics: R², NRMSE, MAE, FIT %, RMSE, Peak Ratio"
    ])

    # SLIDE 12 (HEADLINE RESULTS)
    add_standard_slide("THE HEADLINE: High-Fidelity Convergence", [
        "EMPIRICAL VALIDATION (STAGE A):",
        "• 0.9987 : Voltage R² (TDMS U0 vs model 244·√2·sin(2π·50·t))",
        "• 0.9124 : Current R² (phase-aligned 50 Hz fundamental)",
        "• 99.2 % : Thermal validation (predicted ΔT 11.86 K vs analytical 11.95 K)",
        "• 94.9 % : CR5 mass check (23.74 kg vs 25 kg)",
        "",
        "EXPORT FIDELITY (STAGE B):",
        "• > 99.99 % agreement on all outputs.",
        "• T_motor variance: 0.0001 K",
        "• P_elec variance: 0.0003 W"
    ], (7, 4, 6, 3, "validation_dynamic_v7.png"))

    # SLIDE 13
    add_standard_slide("SCIENTIFIC HONESTY: Quantifying Residuals", [
        "IDENTIFIED RESIDUAL A: 14% Current Amplitude Offset",
        "• Model is grid-connected; real drive has unmodelled rectifier conduction losses.",
        "",
        "IDENTIFIED RESIDUAL B: 105° Phase Lag (−5.83 ms)",
        "• Caused by the exact same unmodelled rectifier component.",
        "",
        "DEFENCE POSITION:",
        "\"Every residual is localised, named, and physically explained.\"",
        "",
        "THE FIX (PHASE 2):",
        "• Parameter calibration in Dymola 2026x Refresh 1 to reduce residual to 3–5 %."
    ], (8, 2, 5, 4, "Zoomed Phase Mismatch Plot"))

    # SLIDE 14
    add_standard_slide("ENERGY ANALYSIS", [
        "TOTAL LINE ENERGY: 3645 J per cycle (~121 W average)",
        "",
        "BREAKDOWN:",
        "• Conveyor: 90.1 % (3284 J)",
        "• Sensors + PLC: 6.6 % (240 J)",
        "• Pneumatic System: 3.3 % (121 J)",
        "• XY10 Steppers: ≈ 0 % (time-averaged DC-equivalent bursts)",
        "",
        "INSIGHT:",
        "• Conveyor is the dominant predictive-maintenance target.",
        "• Model exposes T_motor directly, aligning with Industry 5.0 sustainability KPIs."
    ], (8, 2, 5, 4, "energy_analysis_FIXED.png (Doughnut)"))

    # SLIDE 15
    add_standard_slide("INTEGRATION: Spatial and Coupling", [
        "SPATIAL LAYOUT:",
        "• 8445 mm dimensioned scale drawing (X-Y and X-Z views).",
        "• Explicit tracking of modelled vs unmodelled peripheral zones.",
        "",
        "COUPLING MATRIX:",
        "• 25 interactions mapped across 8 subsystems.",
        "• Types tracked: Electrical, Mechanical, Thermal, EMI, Part-Flow.",
        "• Highlights explicit strong couplings (e.g., Pneumatic ↔ Gripper)."
    ], (6, 3, 7, 4, "alix_spatial_layout.png + coupling_matrix.png"))

    # ==========================================
    # SLIDE 16: SCORECARD (TABLE)
    # ==========================================
    slide16 = prs.slides.add_slide(layout_blank)
    title = slide16.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(10), Inches(1))
    title.text_frame.text = "REQUIREMENTS SCORECARD"
    set_text_format(title.text_frame.paragraphs[0].runs[0], 32, bold=True, color=NAVY)

    rows, cols = 13, 2
    table_shape = slide16.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(11), Inches(5))
    table = table_shape.table

    table.cell(0, 0).text = "Requirement"
    table.cell(0, 1).text = "Status"

    reqs = [
        ("Mechanical modelling", "✓ DONE"),
        ("Electrical modelling", "✓ DONE"),
        ("Thermal modelling", "✓ DONE"),
        ("Spatial representation", "✓ DONE"),
        ("Coupling analysis", "✓ DONE"),
        ("FMU export", "✓ DONE"),
        ("Python orchestration", "✓ DONE"),
        ("End-to-end simulation", "✓ DONE"),
        ("Validation vs measurements", "✓ DONE"),
        ("Energy analysis", "✓ DONE"),
        ("Sensitivity / observability", "⚠ FOUNDATION (Phase 2)"),
        ("Parameter identification", "⚠ FOUNDATION (Phase 2)")
    ]

    for r in range(12):
        table.cell(r + 1, 0).text = reqs[r][0]
        table.cell(r + 1, 1).text = reqs[r][1]

    txBox = slide16.shapes.add_textbox(Inches(1), Inches(6.8), Inches(10), Inches(0.5))
    txBox.text_frame.text = "SUMMARY: 10 of 12 fully delivered · 2 with foundation ready for Phase 2."

    # SLIDE 17
    add_standard_slide("ROADMAP: Model → Shadow → Twin", [
        "PHASE 1: Digital Model [NOW]",
        "• Validated multiphysics model running end-to-end.",
        "",
        "PHASE 2: Calibration",
        "• Tuning in Dymola to cut current residual to 3–5%.",
        "",
        "PHASE 3: Sensitivity & DoE",
        "• Formal design of experiments.",
        "",
        "PHASE 4: Digital Shadow",
        "• Live one-way data link (real-time TDMS streaming).",
        "",
        "PHASE 5: Digital Twin",
        "• Close the loop (reinforcement learning & optimization)."
    ])

    # SLIDE 18
    add_standard_slide("CLOSING SUMMARY", [
        "\"Three subsystems · three physics domains · validated to R² = 0.999 voltage,",
        "0.91 current, 99.2 % thermal. Every residual is physically explained.\"",
        "",
        "CORE TAKEAWAYS:",
        "✓ Validated, not asserted.",
        "✓ Honest by construction.",
        "✓ Built to extend.",
        "",
        "Thank you. Questions welcome.",
        "",
        "Samiksha Satre | Euler Laboratory, ISAE-Supméca",
        "github.com/samikshasatre/Mechatronics-modelling-of-entire-industry-5.0"
    ], (10, 6, 2, 1, "Dassault Systèmes Logo"))

    # SAVE
    prs.save('ALIX_Digital_Twin_Dassault.pptx')
    print("Presentation successfully generated: ALIX_Digital_Twin_Dassault.pptx")


if __name__ == '__main__':
    create_deck()